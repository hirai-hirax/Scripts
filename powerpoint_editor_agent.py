import os
import glob
from pathlib import Path
from typing import List, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches
from pptx.shapes.base import BaseShape
try:
    import comtypes.client
except ImportError:
    comtypes = None


class PowerPointEditorAgent:
    def __init__(self, name="PowerPoint表紙編集エージェント"):
        self.name = name
        self.supported_extensions = ['.pptx', '.ppt']
        self.circled_numbers = {
            '①': '1', '②': '2', '③': '3', '④': '4', '⑤': '5',
            '⑥': '6', '⑦': '7', '⑧': '8', '⑨': '9', '⑩': '10'
        }
    
    def find_powerpoint_files(self, folder_path: str) -> List[str]:
        """指定フォルダ内のPowerPointファイルを検索する"""
        if not os.path.exists(folder_path):
            raise FileNotFoundError(f"フォルダが見つかりません: {folder_path}")
        
        ppt_files = []
        
        for ext in self.supported_extensions:
            pattern = os.path.join(folder_path, f"**/*{ext}")
            files = glob.glob(pattern, recursive=True)
            ppt_files.extend(files)
        
        return sorted(ppt_files)
    
    
    
    def detect_circled_number(self, ppt_file_path: str) -> Optional[str]:
        """PowerPointファイルの表紙から丸付き番号を検出"""
        try:
            prs = Presentation(ppt_file_path)
            if len(prs.slides) == 0:
                return None
            
            first_slide = prs.slides[0]
            
            # すべてのテキストを検索
            for shape in first_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text
                    for circled_num in self.circled_numbers.keys():
                        if circled_num in text:
                            return circled_num
            
            return None
            
        except Exception:
            return None
    
    def find_matching_textbox_in_stamp_file(self, stamp_file_path: str, target_number: str) -> Optional[BaseShape]:
        """スタンプ集.pptxから指定した丸付き番号を含むテキストボックスを検索"""
        try:
            if not os.path.exists(stamp_file_path):
                return None
            
            prs = Presentation(stamp_file_path)
            if len(prs.slides) == 0:
                return None
            
            first_slide = prs.slides[0]
            
            # 指定した丸付き番号を含むテキストボックスを検索
            for shape in first_slide.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    text = shape.text_frame.text
                    if target_number in text:
                        return shape
            
            return None
            
        except Exception:
            return None
    
    def copy_textbox_to_slide(self, source_shape: BaseShape, target_slide, position: Tuple[float, float] = None):
        """テキストボックスを指定位置にコピー"""
        try:
            # デフォルト位置は左上
            if position is None:
                left = Inches(0.5)
                top = Inches(0.5)
            else:
                left = Inches(position[0])
                top = Inches(position[1])
            
            # 元の形状のサイズを取得
            width = source_shape.width
            height = source_shape.height
            
            # 新しいテキストボックスを作成
            new_textbox = target_slide.shapes.add_textbox(left, top, width, height)
            
            # テキスト内容をコピー
            if hasattr(source_shape, 'text_frame') and source_shape.text_frame:
                source_text_frame = source_shape.text_frame
                target_text_frame = new_textbox.text_frame
                
                # 既存のパラグラフをクリア
                target_text_frame.clear()
                
                # パラグラフを一つずつコピー
                for i, source_paragraph in enumerate(source_text_frame.paragraphs):
                    if i == 0:
                        target_paragraph = target_text_frame.paragraphs[0]
                    else:
                        target_paragraph = target_text_frame.add_paragraph()
                    
                    target_paragraph.text = source_paragraph.text
                    target_paragraph.alignment = source_paragraph.alignment
                    
                    # フォント設定をコピー
                    if source_paragraph.font:
                        target_paragraph.font.name = source_paragraph.font.name
                        target_paragraph.font.size = source_paragraph.font.size
                        if source_paragraph.font.color.rgb:
                            target_paragraph.font.color.rgb = source_paragraph.font.color.rgb
            
            return new_textbox
            
        except Exception as e:
            print(f"テキストボックスのコピーに失敗: {str(e)}")
            return None
    
    def add_stamp_to_cover_slide(self, ppt_file_path: str, stamp_file_path: str, 
                                position: Tuple[float, float] = None, save_as: Optional[str] = None) -> str:
        """PowerPointの表紙に丸付き番号に対応するスタンプを追加"""
        try:
            # 丸付き番号を検出
            circled_number = self.detect_circled_number(ppt_file_path)
            if not circled_number:
                raise ValueError("表紙に丸付き番号が見つかりませんでした")
            
            print(f"検出された丸付き番号: {circled_number}")
            
            # スタンプファイルから対応するテキストボックスを検索
            matching_textbox = self.find_matching_textbox_in_stamp_file(stamp_file_path, circled_number)
            if not matching_textbox:
                raise ValueError(f"スタンプファイルに'{circled_number}'を含むテキストボックスが見つかりませんでした")
            
            # 対象PowerPointを開く
            prs = Presentation(ppt_file_path)
            if len(prs.slides) == 0:
                raise ValueError("プレゼンテーションにスライドが含まれていません")
            
            first_slide = prs.slides[0]
            
            # テキストボックスをコピー
            copied_textbox = self.copy_textbox_to_slide(matching_textbox, first_slide, position)
            if not copied_textbox:
                raise ValueError("テキストボックスのコピーに失敗しました")
            
            # 保存
            if save_as:
                output_path = save_as
            else:
                file_path = Path(ppt_file_path)
                output_path = str(file_path.parent / f"{file_path.stem}_stamped{file_path.suffix}")
            
            prs.save(output_path)
            return output_path
            
        except Exception as e:
            raise Exception(f"スタンプ追加に失敗しました: {str(e)}")
    
    def batch_add_stamps_to_files(self, folder_path: str, stamp_file_path: str,
                                 position: Tuple[float, float] = None, output_folder: Optional[str] = None) -> List[str]:
        """フォルダ内の全PowerPointファイルに一括でスタンプを追加"""
        ppt_files = self.find_powerpoint_files(folder_path)
        
        if not ppt_files:
            print(f"フォルダ内にPowerPointファイルが見つかりませんでした: {folder_path}")
            return []
        
        processed_files = []
        
        for ppt_file in ppt_files:
            try:
                print(f"処理中: {os.path.basename(ppt_file)}")
                
                if output_folder:
                    os.makedirs(output_folder, exist_ok=True)
                    file_name = Path(ppt_file).name
                    save_path = os.path.join(output_folder, file_name)
                else:
                    save_path = None
                
                output_path = self.add_stamp_to_cover_slide(ppt_file, stamp_file_path, position, save_path)
                processed_files.append(output_path)
                print(f"完了: {os.path.basename(output_path)}")
                
            except Exception as e:
                print(f"エラー ({os.path.basename(ppt_file)}): {str(e)}")
                continue
        
        return processed_files
    
    def preview_files(self, folder_path: str) -> None:
        """指定フォルダ内のPowerPointファイル一覧を表示"""
        ppt_files = self.find_powerpoint_files(folder_path)
        
        if not ppt_files:
            print(f"フォルダ内にPowerPointファイルが見つかりませんでした: {folder_path}")
            return
        
        print(f"\n見つかったPowerPointファイル ({len(ppt_files)}件):")
        print("-" * 60)
        
        for i, ppt_file in enumerate(ppt_files, 1):
            rel_path = os.path.relpath(ppt_file, folder_path)
            file_size = os.path.getsize(ppt_file) / (1024 * 1024)  # MB
            
            # 丸付き番号を検出
            circled_number = self.detect_circled_number(ppt_file)
            number_info = f" [{circled_number}]" if circled_number else " [番号なし]"
            
            print(f"{i:2d}. {rel_path} ({file_size:.1f}MB){number_info}")
    
    def convert_pptx_to_2up_pdf(self, ppt_file_path: str, output_pdf_path: Optional[str] = None) -> str:
        """指定PowerPointファイルを2アップ縦並びPDFとして出力"""
        if comtypes is None:
            raise ImportError("comtypesライブラリがインストールされていません。pip install comtypesでインストールしてください。")
        
        if not os.path.exists(ppt_file_path):
            raise FileNotFoundError(f"ファイルが見つかりません: {ppt_file_path}")
        
        try:
            # PowerPointアプリケーションを起動
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = False  # 非表示で実行
            
            # ファイルを開く
            abs_path = os.path.abspath(ppt_file_path)
            presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True)
            
            # 出力ファイル名を決定
            if output_pdf_path is None:
                file_path = Path(ppt_file_path)
                output_pdf_path = str(file_path.parent / f"{file_path.stem}_2up.pdf")
            
            abs_output_path = os.path.abspath(output_pdf_path)
            
            # PDFエクスポート設定
            # ppPrintOutputSlides = 1 (スライド)
            # ppPrintHandoutVerticalFirst = 1 (縦方向優先)
            # ppPrintHandout2SlidePerPage = 2 (2スライド/ページ)
            
            # 印刷設定をハンドアウトモードに変更
            presentation.PrintOptions.OutputType = 2  # ppPrintOutputHandouts
            presentation.PrintOptions.HandoutOrder = 1  # ppPrintHandoutVerticalFirst
            presentation.PrintOptions.NumberOfHandoutSlides = 2  # 2スライド/ページ
            
            # PDFとしてエクスポート
            # ppFixedFormatTypePDF = 2
            # ppFixedFormatIntentPrint = 2
            # msoTriStateFalse = 0
            # msoTriStateTrue = -1
            presentation.ExportAsFixedFormat(
                abs_output_path,  # OutputFileName
                2,  # ExportFormat (ppFixedFormatTypePDF)
                2,  # Intent (ppFixedFormatIntentPrint)
                0,  # FrameSlides (msoTriStateFalse)
                2,  # HandoutOrder (ppPrintHandoutVerticalFirst)
                1,  # OutputType (ppPrintOutputHandouts)
                0,  # PrintHiddenSlides (msoTriStateFalse)
                None,  # PrintRange
                1,  # RangeType (ppPrintAll)
                "",  # SlideShowName
                False,  # IncludeDocProps
                True,  # KeepIRMSettings
                True,  # DocStructureTags
                True,  # BitmapMissingFonts
                True,  # UseDocumentICCProfile
                2,  # HandoutSlides (2スライド/ページ)
            )
            
            # プレゼンテーションを閉じる
            presentation.Close()
            
            # PowerPointアプリケーションを終了
            powerpoint.Quit()
            
            return abs_output_path
            
        except Exception as e:
            # エラーが発生した場合もPowerPointを閉じる
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'powerpoint' in locals():
                    powerpoint.Quit()
            except:
                pass
            raise Exception(f"PDF変換に失敗しました: {str(e)}")
    
    def batch_convert_to_2up_pdf(self, folder_path: str, output_folder: Optional[str] = None) -> List[str]:
        """フォルダ内の全PowerPointファイルを2アップPDFに一括変換"""
        ppt_files = self.find_powerpoint_files(folder_path)
        
        if not ppt_files:
            print(f"フォルダ内にPowerPointファイルが見つかりませんでした: {folder_path}")
            return []
        
        if output_folder:
            os.makedirs(output_folder, exist_ok=True)
        
        processed_files = []
        
        for ppt_file in ppt_files:
            try:
                print(f"変換中: {os.path.basename(ppt_file)}")
                
                if output_folder:
                    file_name = Path(ppt_file).stem + "_2up.pdf"
                    output_path = os.path.join(output_folder, file_name)
                else:
                    output_path = None
                
                result_path = self.convert_pptx_to_2up_pdf(ppt_file, output_path)
                processed_files.append(result_path)
                print(f"完了: {os.path.basename(result_path)}")
                
            except Exception as e:
                print(f"エラー ({os.path.basename(ppt_file)}): {str(e)}")
                continue
        
        return processed_files
    
    def convert_pptx_to_pdf(self, ppt_file_path: str, output_pdf_path: Optional[str] = None) -> str:
        """指定PowerPointファイルを通常PDFとして出力"""
        if comtypes is None:
            raise ImportError("comtypesライブラリがインストールされていません。pip install comtypesでインストールしてください。")
        
        if not os.path.exists(ppt_file_path):
            raise FileNotFoundError(f"ファイルが見つかりません: {ppt_file_path}")
        
        try:
            # PowerPointアプリケーションを起動
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = False  # 非表示で実行
            
            # ファイルを開く
            abs_path = os.path.abspath(ppt_file_path)
            presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True)
            
            # 出力ファイル名を決定
            if output_pdf_path is None:
                file_path = Path(ppt_file_path)
                output_pdf_path = str(file_path.parent / f"{file_path.stem}.pdf")
            
            abs_output_path = os.path.abspath(output_pdf_path)
            
            # 通常のスライド出力設定
            presentation.PrintOptions.OutputType = 1  # ppPrintOutputSlides
            
            # PDFとしてエクスポート
            presentation.ExportAsFixedFormat(
                abs_output_path,  # OutputFileName
                2,  # ExportFormat (ppFixedFormatTypePDF)
                2,  # Intent (ppFixedFormatIntentPrint)
                0,  # FrameSlides (msoTriStateFalse)
                1,  # HandoutOrder (ppPrintHandoutVerticalFirst)
                1,  # OutputType (ppPrintOutputSlides)
                0,  # PrintHiddenSlides (msoTriStateFalse)
                None,  # PrintRange
                1,  # RangeType (ppPrintAll)
                "",  # SlideShowName
                False,  # IncludeDocProps
                True,  # KeepIRMSettings
                True,  # DocStructureTags
                True,  # BitmapMissingFonts
                True,  # UseDocumentICCProfile
            )
            
            # プレゼンテーションを閉じる
            presentation.Close()
            
            # PowerPointアプリケーションを終了
            powerpoint.Quit()
            
            return abs_output_path
            
        except Exception as e:
            # エラーが発生した場合もPowerPointを閉じる
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'powerpoint' in locals():
                    powerpoint.Quit()
            except:
                pass
            raise Exception(f"PDF変換に失敗しました: {str(e)}")
    
    def batch_convert_to_pdf(self, folder_path: str, output_folder: Optional[str] = None) -> List[str]:
        """フォルダ内の全PowerPointファイルを通常PDFに一括変換"""
        ppt_files = self.find_powerpoint_files(folder_path)
        
        if not ppt_files:
            print(f"フォルダ内にPowerPointファイルが見つかりませんでした: {folder_path}")
            return []
        
        if output_folder:
            os.makedirs(output_folder, exist_ok=True)
        
        processed_files = []
        
        for ppt_file in ppt_files:
            try:
                print(f"変換中: {os.path.basename(ppt_file)}")
                
                if output_folder:
                    file_name = Path(ppt_file).stem + ".pdf"
                    output_path = os.path.join(output_folder, file_name)
                else:
                    output_path = None
                
                result_path = self.convert_pptx_to_pdf(ppt_file, output_path)
                processed_files.append(result_path)
                print(f"完了: {os.path.basename(result_path)}")
                
            except Exception as e:
                print(f"エラー ({os.path.basename(ppt_file)}): {str(e)}")
                continue
        
        return processed_files
    
    def print_pptx_2up(self, ppt_file_path: str, printer_name: Optional[str] = None, copies: int = 1) -> bool:
        """指定PowerPointファイルを2アップ縦方向でプリンタ印刷"""
        if comtypes is None:
            raise ImportError("comtypesライブラリがインストールされていません。pip install comtypesでインストールしてください。")
        
        if not os.path.exists(ppt_file_path):
            raise FileNotFoundError(f"ファイルが見つかりません: {ppt_file_path}")
        
        try:
            # PowerPointアプリケーションを起動
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = False  # 非表示で実行
            
            # ファイルを開く
            abs_path = os.path.abspath(ppt_file_path)
            presentation = powerpoint.Presentations.Open(abs_path, ReadOnly=True)
            
            # 印刷設定を2アップハンドアウトモードに変更
            presentation.PrintOptions.OutputType = 2  # ppPrintOutputHandouts
            presentation.PrintOptions.HandoutOrder = 1  # ppPrintHandoutVerticalFirst (縦方向優先)
            presentation.PrintOptions.NumberOfHandoutSlides = 2  # 2スライド/ページ
            presentation.PrintOptions.Copies = copies
            
            # プリンタを指定（指定がない場合はデフォルト）
            if printer_name:
                presentation.PrintOptions.ActivePrinter = printer_name
            
            # 印刷範囲を全スライドに設定
            presentation.PrintOptions.RangeType = 1  # ppPrintAll
            
            # 印刷実行
            presentation.PrintOut()
            
            print(f"印刷ジョブを送信しました: {os.path.basename(ppt_file_path)}")
            
            # プレゼンテーションを閉じる
            presentation.Close()
            
            # PowerPointアプリケーションを終了
            powerpoint.Quit()
            
            return True
            
        except Exception as e:
            # エラーが発生した場合もPowerPointを閉じる
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'powerpoint' in locals():
                    powerpoint.Quit()
            except:
                pass
            raise Exception(f"印刷に失敗しました: {str(e)}")
    
    def batch_print_2up(self, folder_path: str, printer_name: Optional[str] = None, copies: int = 1) -> List[str]:
        """フォルダ内の全PowerPointファイルを2アップで一括印刷"""
        ppt_files = self.find_powerpoint_files(folder_path)
        
        if not ppt_files:
            print(f"フォルダ内にPowerPointファイルが見つかりませんでした: {folder_path}")
            return []
        
        printed_files = []
        
        for ppt_file in ppt_files:
            try:
                print(f"印刷中: {os.path.basename(ppt_file)}")
                
                success = self.print_pptx_2up(ppt_file, printer_name, copies)
                if success:
                    printed_files.append(ppt_file)
                    print(f"印刷完了: {os.path.basename(ppt_file)}")
                
            except Exception as e:
                print(f"エラー ({os.path.basename(ppt_file)}): {str(e)}")
                continue
        
        return printed_files
    
    def get_available_printers(self) -> List[str]:
        """利用可能なプリンタ一覧を取得"""
        if comtypes is None:
            return []
        
        try:
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = False
            
            # 新しい空のプレゼンテーションを作成
            presentation = powerpoint.Presentations.Add()
            
            # インストールされているプリンタの取得は複雑なので、
            # 代わりにWindowsのプリンタ情報を取得
            import win32print
            printers = [printer[2] for printer in win32print.EnumPrinters(2)]
            
            presentation.Close()
            powerpoint.Quit()
            
            return printers
            
        except Exception:
            # win32printが利用できない場合やエラーの場合は空リストを返す
            try:
                if 'presentation' in locals():
                    presentation.Close()
                if 'powerpoint' in locals():
                    powerpoint.Quit()
            except:
                pass
            return []


def main():
    agent = PowerPointEditorAgent()
    
    print(f"\n{agent.name}が起動しました！")
    print("=" * 50)
    
    while True:
        print("\n操作を選択してください:")
        print("1. フォルダ内のPowerPointファイル一覧表示")
        print("2. 単一ファイルにスタンプ追加")
        print("3. フォルダ内全ファイルに一括スタンプ追加")
        print("4. 単一ファイルを通常PDF変換")
        print("5. フォルダ内全ファイルを通常PDF一括変換")
        print("6. 単一ファイルを2アップPDF変換")
        print("7. フォルダ内全ファイルを2アップPDF一括変換")
        print("8. 単一ファイルを2アッププリンタ印刷")
        print("9. フォルダ内全ファイルを2アップ一括印刷")
        print("10. 終了")
        
        choice = input("\n選択 (1-10): ").strip()
        
        if choice == "1":
            folder_path = input("フォルダパスを入力: ").strip()
            if folder_path:
                try:
                    agent.preview_files(folder_path)
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "2":
            ppt_file = input("PowerPointファイルのパスを入力: ").strip()
            stamp_file = input("スタンプ集.pptxのパスを入力: ").strip()
            
            if ppt_file and stamp_file:
                try:
                    output_path = agent.add_stamp_to_cover_slide(ppt_file, stamp_file)
                    print(f"完了！保存先: {output_path}")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "3":
            folder_path = input("フォルダパスを入力: ").strip()
            stamp_file = input("スタンプ集.pptxのパスを入力: ").strip()
            output_folder = input("出力フォルダ (空白で元ファイルを上書き): ").strip() or None
            
            if folder_path and stamp_file:
                try:
                    processed_files = agent.batch_add_stamps_to_files(folder_path, stamp_file, None, output_folder)
                    print(f"\n処理完了！ {len(processed_files)}件のファイルを処理しました。")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "4":
            ppt_file = input("PowerPointファイルのパスを入力: ").strip()
            
            if ppt_file:
                try:
                    output_path = agent.convert_pptx_to_pdf(ppt_file)
                    print(f"完了！保存先: {output_path}")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "5":
            folder_path = input("フォルダパスを入力: ").strip()
            output_folder = input("出力PDFフォルダ (空白で元フォルダに保存): ").strip() or None
            
            if folder_path:
                try:
                    processed_files = agent.batch_convert_to_pdf(folder_path, output_folder)
                    print(f"\n変換完了！ {len(processed_files)}件のPDFファイルを作成しました。")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "6":
            ppt_file = input("PowerPointファイルのパスを入力: ").strip()
            
            if ppt_file:
                try:
                    output_path = agent.convert_pptx_to_2up_pdf(ppt_file)
                    print(f"完了！保存先: {output_path}")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "7":
            folder_path = input("フォルダパスを入力: ").strip()
            output_folder = input("出力PDFフォルダ (空白で元フォルダに保存): ").strip() or None
            
            if folder_path:
                try:
                    processed_files = agent.batch_convert_to_2up_pdf(folder_path, output_folder)
                    print(f"\n変換完了！ {len(processed_files)}件のPDFファイルを作成しました。")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "8":
            ppt_file = input("PowerPointファイルのパスを入力: ").strip()
            
            if ppt_file:
                # 利用可能なプリンタ一覧を表示
                printers = agent.get_available_printers()
                if printers:
                    print("\n利用可能なプリンタ:")
                    for i, printer in enumerate(printers, 1):
                        print(f"{i}. {printer}")
                    print("0. デフォルトプリンタ")
                    
                    printer_choice = input("\nプリンタを選択 (0-{}): ".format(len(printers))).strip()
                    if printer_choice.isdigit():
                        choice_num = int(printer_choice)
                        if choice_num == 0:
                            selected_printer = None
                        elif 1 <= choice_num <= len(printers):
                            selected_printer = printers[choice_num - 1]
                        else:
                            selected_printer = None
                    else:
                        selected_printer = None
                else:
                    selected_printer = None
                
                copies_input = input("印刷部数 [1]: ").strip()
                copies = int(copies_input) if copies_input.isdigit() and int(copies_input) > 0 else 1
                
                try:
                    success = agent.print_pptx_2up(ppt_file, selected_printer, copies)
                    if success:
                        print(f"印刷ジョブを送信しました！")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "9":
            folder_path = input("フォルダパスを入力: ").strip()
            
            if folder_path:
                # 利用可能なプリンタ一覧を表示
                printers = agent.get_available_printers()
                if printers:
                    print("\n利用可能なプリンタ:")
                    for i, printer in enumerate(printers, 1):
                        print(f"{i}. {printer}")
                    print("0. デフォルトプリンタ")
                    
                    printer_choice = input("\nプリンタを選択 (0-{}): ".format(len(printers))).strip()
                    if printer_choice.isdigit():
                        choice_num = int(printer_choice)
                        if choice_num == 0:
                            selected_printer = None
                        elif 1 <= choice_num <= len(printers):
                            selected_printer = printers[choice_num - 1]
                        else:
                            selected_printer = None
                    else:
                        selected_printer = None
                else:
                    selected_printer = None
                
                copies_input = input("印刷部数 [1]: ").strip()
                copies = int(copies_input) if copies_input.isdigit() and int(copies_input) > 0 else 1
                
                try:
                    printed_files = agent.batch_print_2up(folder_path, selected_printer, copies)
                    print(f"\n印刷完了！ {len(printed_files)}件のファイルを印刷しました。")
                except Exception as e:
                    print(f"エラー: {e}")
        
        elif choice == "10":
            break
        
        else:
            print("無効な選択です。1-10を入力してください。")
    
    print(f"\n{agent.name}をご利用いただき、ありがとうございました！")


if __name__ == "__main__":
    main()