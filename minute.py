import os
import streamlit as st
import pandas as pd
from openai import AzureOpenAI
from io import BytesIO
from dotenv import load_dotenv
from pptx import Presentation
import docx # For .docx files
from docx import Document # Explicitly import Document for creation
import fitz # PyMuPDF for .pdf files

# .envファイルから環境変数を読み込む
load_dotenv()

# 環境変数から設定を取得（Azure OpenAI のエンドポイント・API キーを設定してください）
AZURE_OPENAI_ENDPOINT = os.environ.get("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.environ.get("AZURE_OPENAI_API_KEY")
API_VERSION = "2024-02-01" # ご利用の API バージョンに合わせてください

@st.cache_resource
def get_openai_client():
    if not AZURE_OPENAI_API_KEY or not AZURE_OPENAI_ENDPOINT:
        st.error("Azure OpenAI APIキーまたはエンドポイントが設定されていません。環境変数を確認してください。")
        return None
    return AzureOpenAI(
        api_key=AZURE_OPENAI_API_KEY,
        api_version=API_VERSION,
        azure_endpoint=AZURE_OPENAI_ENDPOINT
    )

# --- Text Extraction Functions ---
def extract_text_from_pptx(file: BytesIO):
    """Extracts all text from a PowerPoint (PPTX) file."""
    try:
        prs = Presentation(file)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        st.warning(f"PowerPointファイルの読み込み中にエラーが発生しました: {e}")
        return None

def extract_text_from_docx(file: BytesIO):
    """Extracts all text from a Word (DOCX) file."""
    try:
        doc = docx.Document(file)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except Exception as e:
        st.warning(f"Wordファイルの読み込み中にエラーが発生しました: {e}")
        return None

def extract_text_from_pdf(file: BytesIO):
    """Extracts all text from a PDF file."""
    try:
        # file.read() is needed for fitz.open(stream=...)
        file_bytes = file.read()
        pdf_document = fitz.open(stream=file_bytes, filetype='pdf')
        text = ""
        for page in pdf_document:
            text += page.get_text()
        pdf_document.close()
        return text
    except Exception as e:
        st.warning(f"PDFファイルの読み込み中にエラーが発生しました: {e}")
        return None

def extract_text_from_excel_for_agenda(file: BytesIO):
    """
    Reads an Excel file and extracts text content row by row from all sheets for agenda generation.
    Combines all non-empty string values from each row.
    """
    try:
        xls = pd.ExcelFile(file)
        full_text_lines = []
        for sheet_name in xls.sheet_names:
            df = xls.parse(sheet_name)
            for index, row in df.iterrows():
                row_text_parts = []
                for col in df.columns:
                    cell_value = row[col]
                    if pd.notna(cell_value) and str(cell_value).strip() != "":
                        row_text_parts.append(str(cell_value).strip())
                if row_text_parts:
                    full_text_lines.append(" ".join(row_text_parts))
        return "\n".join(full_text_lines)
    except Exception as e:
        st.warning(f"Excelファイルの読み込み中にエラーが発生しました: {e}")
        return None

def extract_text_from_excel_for_minutes(file: BytesIO):
    """
    Reads an Excel file and extracts text content from 'speaker' and 'text' columns,
    combining them row by row for meeting minutes generation.
    """
    try:
        df = pd.read_excel(file)
        
        if 'speaker' not in df.columns and 'text' not in df.columns:
            st.error("Excelファイルに'speaker'列または'text'列が見つかりません。")
            return None
        
        transcription_lines = []
        for index, row in df.iterrows():
            speaker = row['speaker'] if 'speaker' in df.columns and pd.notna(row['speaker']) else ""
            text = row['text'] if 'text' in df.columns and pd.notna(row['text']) else ""
            
            if speaker and text:
                transcription_lines.append(f"{speaker}: {text}")
            elif text: # If only text is present
                transcription_lines.append(text)
                
        return "\n".join(transcription_lines)
    except Exception as e:
        st.error(f"Excelファイルの読み込み中にエラーが発生しました: {e}")
        return None

# --- Prompt Templates ---
MEETING_MINUTES_PROMPT_TEMPLATE = """
以下の会議の文字起こしデータから、議事録を作成してください。
議事録は、以下の形式で通常の文章で作成してください。箇条書きを使用する場合は、ハイフン(-)やアスタリスク(*)などのマークダウン記号は使用せず、全角の「・」を使用してください。

1. 会議の目的
    ・（会議の目的を簡潔に記載）

2. 参加者（もし文字起こしから推測できる場合）
    ・○○部: ～～さん、△△さん
    ・□□部: ～～さん、△△さん
    …（必要に応じて追加）

3. 主要なトピックと決定事項
    (1)○○について：
        ～～（議論のポイントを記載）
    (2)△△について：
        ～～（議論のポイントを記載）
    …（必要に応じて追加）

4. 今後のアクションアイテムと担当者（もし文字起こしから推測できる場合）
    ・○○について、～～～を行う（担当:△△さん）
    …（必要に応じて追加）
    
5. 次回会議の予定（もし文字起こしから推測できる場合）

---
文字起こしデータ:
{transcription_text}
---

議事録:
"""

MEETING_AGENDA_PROMPT_TEMPLATE = """
以下の提供された資料の内容に基づいて、次回の会議のアジェンダを作成してください。
複数のファイルが提供される場合、それらの内容を統合してアジェンダを作成してください。
アジェンダは、以下の形式で通常の文章で作成してください。箇条書きを使用する場合は、ハイフン(-)やアスタリスク(*)などのマークダウン記号は使用せず、全角の「・」を使用してください。

1. 会議の目的
    ・（会議の目的を簡潔に記載）

2. 主要な議題とそれぞれの簡単な説明
    (1)○○について
    (2)△△について
    …（必要に応じて追加）

3. 議論すべきポイントや決定事項の候補
    (1)○○について：
        ～～（議論のポイントを記載）
    (2)△△について：
        ～～（議論のポイントを記載）
    …（必要に応じて追加）


4. 参加者（もし資料から推測できる場合）
    ・○○部: ～～さん、△△さん
    …（必要に応じて追加）

5. 必要な準備事項（もし資料から推測できる場合）
    ・○○について、～～～を行う（担当:△△さん）
    …（必要に応じて追加）

---
資料の内容:
{combined_presentation_text}
---

会議アジェンダ:
"""

# --- Generation Functions ---
def generate_meeting_minutes_from_text(client: AzureOpenAI, transcription_text: str):
    """
    Generates meeting minutes from a given transcription text using Azure OpenAI Service.
    """
    if client is None:
        return "エラー: Azure OpenAIクライアントが初期化されていません。"

    try:
        prompt_with_text = MEETING_MINUTES_PROMPT_TEMPLATE.format(transcription_text=transcription_text)

        st.info("Azure OpenAI Serviceで議事録を生成中...")
        response = client.chat.completions.create(
            model="gpt-4o",  # Replace with your deployed model name
            messages=[
                {"role": "system", "content": "あなたは会議の文字起こしデータから議事録を作成するアシスタントです。"},
                {"role": "user", "content": prompt_with_text}
            ],
            temperature=0.7,
            max_tokens=1500
        )

        meeting_minutes = response.choices[0].message.content
        st.success("議事録の生成が完了しました！")
        return meeting_minutes

    except Exception as e:
        st.error(f"議事録生成中にエラーが発生しました: {e}")
        return f"エラー: {e}"

def generate_meeting_agenda_from_text(client: AzureOpenAI, combined_presentation_text: str):
    """
    Generates a meeting agenda from given combined presentation text using Azure OpenAI Service.
    """
    if client is None:
        return "エラー: Azure OpenAIクライアントが初期化されていません。"

    try:
        prompt_with_text = MEETING_AGENDA_PROMPT_TEMPLATE.format(combined_presentation_text=combined_presentation_text)

        st.info("Azure OpenAI Serviceで会議アジェンダを生成中...")
        response = client.chat.completions.create(
            model="gpt-4o",  # Replace with your deployed model name
            messages=[
                {"role": "system", "content": "あなたは提供された資料の内容から会議アジェンダを作成するアシスタントです。"},
                {"role": "user", "content": prompt_with_text}
            ],
            temperature=0.7,
            max_tokens=1500
        )

        meeting_agenda = response.choices[0].message.content
        st.success("会議アジェンダの生成が完了しました！")
        return meeting_agenda

    except Exception as e:
        st.error(f"会議アジェンダ生成中にエラーが発生しました: {e}")
        return f"エラー: {e}"

# --- Streamlit UI Functions ---
def create_docx_from_text(text: str) -> BytesIO:
    """Creates a BytesIO object containing a DOCX file from a given text."""
    document = Document()
    # Split text by lines and add as paragraphs
    for line in text.split('\n'):
        document.add_paragraph(line)
    
    buffer = BytesIO()
    document.save(buffer)
    buffer.seek(0)
    return buffer

def meeting_minutes_app():
    st.title("文字起こしデータから議事録作成アプリ")
    st.write("文字起こしファイルをアップロードすると、Azure OpenAI Serviceを使って議事録を自動生成します。")

    uploaded_file = st.file_uploader(
        "文字起こしファイルを選択してください (TXTまたはExcel)",
        type=["txt", "xlsx"],
        key="transcription_file_uploader"
    )

    transcription_text = None
    if uploaded_file is not None:
        file_extension = uploaded_file.name.split('.')[-1].lower()

        if file_extension == "txt":
            transcription_text = uploaded_file.read().decode('utf-8')
            st.subheader("アップロードされた文字起こしデータ (TXT)")
            st.text_area("プレビュー", transcription_text, height=200)
        elif file_extension == "xlsx":
            transcription_text = extract_text_from_excel_for_minutes(uploaded_file)
            if transcription_text:
                st.subheader("アップロードされた文字起こしデータ (Excelから抽出)")
                st.text_area("プレビュー", transcription_text, height=200)
        else:
            st.error("サポートされていないファイル形式です。TXTまたはExcelファイルをアップロードしてください。")

    if st.button("議事録を生成", key="generate_minutes_button"):
        if transcription_text:
            client = get_openai_client()
            if client:
                with st.spinner("議事録を生成中..."):
                    meeting_minutes = generate_meeting_minutes_from_text(client, transcription_text)
                    
                    st.subheader("生成された議事録")
                    st.text_area("議事録", meeting_minutes, height=400)

                if meeting_minutes and not meeting_minutes.startswith("エラー:"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.download_button(
                            label="議事録をダウンロード (TXT)",
                            data=meeting_minutes.encode('utf-8'),
                            file_name="generated_meeting_minutes.txt",
                            mime="text/plain"
                        )
                    with col2:
                        docx_buffer = create_docx_from_text(meeting_minutes)
                        st.download_button(
                            label="議事録をダウンロード (DOCX)",
                            data=docx_buffer,
                            file_name="generated_meeting_minutes.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )
            else:
                meeting_minutes = "エラー: Azure OpenAIクライアントが初期化されていません。"
                st.text_area("議事録", meeting_minutes, height=400)
        else:
            st.warning("文字起こしファイルをアップロードしてください。")

def meeting_agenda_app():
    st.title("資料から会議アジェンダ作成アプリ")
    st.write("PowerPoint、Word、TXT、Excel、PDFファイルを複数アップロードすると、Azure OpenAI Serviceを使って次回の会議アジェンダを自動生成します。")

    uploaded_files = st.file_uploader(
        "資料ファイルを選択してください (PPTX, DOCX, TXT, XLSX, PDF)",
        type=["pptx", "docx", "txt", "xlsx", "pdf"],
        accept_multiple_files=True,
        key="presentation_files_uploader"
    )

    combined_presentation_text = ""
    if uploaded_files:
        st.subheader("アップロードされた資料の内容プレビュー")
        for uploaded_file in uploaded_files:
            file_extension = uploaded_file.name.split('.')[-1].lower()
            file_content = None

            if file_extension == "pptx":
                file_content = extract_text_from_pptx(uploaded_file)
            elif file_extension == "docx":
                file_content = extract_text_from_docx(uploaded_file)
            elif file_extension == "txt":
                file_content = uploaded_file.read().decode('utf-8')
            elif file_extension == "xlsx":
                file_content = extract_text_from_excel_for_agenda(uploaded_file)
            elif file_extension == "pdf":
                file_content = extract_text_from_pdf(uploaded_file)
            else:
                st.warning(f"サポートされていないファイル形式です: {uploaded_file.name}")
                continue

            if file_content:
                combined_presentation_text += f"\n--- {uploaded_file.name} ---\n"
                combined_presentation_text += file_content
                st.text_area(f"ファイル: {uploaded_file.name}", file_content, height=150, key=f"preview_{uploaded_file.name}")
            else:
                st.warning(f"ファイル '{uploaded_file.name}' からテキストを抽出できませんでした。")

    if st.button("アジェンダを生成", key="generate_agenda_button"):
        if combined_presentation_text:
            client = get_openai_client()
            if client:
                with st.spinner("アジェンダを生成中..."):
                    meeting_agenda = generate_meeting_agenda_from_text(client, combined_presentation_text)
                    
                    st.subheader("生成された会議アジェンダ")
                    st.text_area("会議アジェンダ", meeting_agenda, height=400)

                if meeting_agenda and not meeting_agenda.startswith("エラー:"):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.download_button(
                            label="アジェンダをダウンロード (TXT)",
                            data=meeting_agenda.encode('utf-8'),
                            file_name="generated_meeting_agenda.txt",
                            mime="text/plain"
                        )
                    with col2:
                        docx_buffer = create_docx_from_text(meeting_agenda)
                        st.download_button(
                            label="アジェンダをダウンロード (DOCX)",
                            data=docx_buffer,
                            file_name="generated_meeting_agenda.docx",
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        )
            else:
                meeting_agenda = "エラー: Azure OpenAIクライアントが初期化されていません。"
                st.text_area("会議アジェンダ", meeting_agenda, height=400)
        else:
            st.warning("資料ファイルをアップロードしてください。")

def main():
    st.set_page_config(layout="wide")
    st.sidebar.title("アプリ選択")
    app_mode = st.sidebar.selectbox(
        "実行するアプリを選択してください:",
        ["議事録作成アプリ", "アジェンダ作成アプリ"]
    )

    if app_mode == "議事録作成アプリ":
        meeting_minutes_app()
    elif app_mode == "アジェンダ作成アプリ":
        meeting_agenda_app()

if __name__ == "__main__":
    main()
