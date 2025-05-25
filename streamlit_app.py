#!/usr/bin/env python3
import os
from pathlib import Path
import streamlit as st
from dotenv import load_dotenv
import pandas as pd
import fitz
import docx
from pptx import Presentation
from langchain_openai import AzureOpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.docstore.document import Document
from openai import AzureOpenAI

# Load environment variables
load_dotenv()
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_API_VERSION = os.getenv("AZURE_OPENAI_API_VERSION")
EMBEDDING_DEPLOYMENT = os.getenv("AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME")
CHAT_DEPLOYMENT = os.getenv("AZURE_OPENAI_CHAT_DEPLOYMENT_NAME")

if not all([AZURE_OPENAI_API_KEY, AZURE_OPENAI_ENDPOINT, EMBEDDING_DEPLOYMENT, CHAT_DEPLOYMENT]):
    st.error(
        "Ensure AZURE_OPENAI_KEY, AZURE_OPENAI_ENDPOINT, "
        "AZURE_OPENAI_EMBEDDING_DEPLOYMENT_NAME, and AZURE_OPENAI_CHAT_DEPLOYMENT_NAME are set in .env"
    )
    st.stop()

# Initialize Azure OpenAI client for chat
client = AzureOpenAI(
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    api_key=AZURE_OPENAI_API_KEY,
    api_version=AZURE_API_VERSION
)

# Initialize LangChain embeddings for Azure OpenAI
embeddings = AzureOpenAIEmbeddings(
    model=EMBEDDING_DEPLOYMENT,
    azure_endpoint=AZURE_OPENAI_ENDPOINT,
    openai_api_key=AZURE_OPENAI_API_KEY,
    openai_api_version=AZURE_API_VERSION,
)

# Directories and paths
NOTES_DIR = Path(__file__).parent / "notes"
NOTES_DIR.mkdir(exist_ok=True)
STORE_DIR = NOTES_DIR / "faiss_store"
STORE_DIR.mkdir(exist_ok=True)

# Directory for converted markdown files
MD_DIR = NOTES_DIR / "markdowns"
MD_DIR.mkdir(exist_ok=True)

def slugify(title: str) -> str:
    safe = "".join(c for c in title if c.isalnum() or c in (" ", "-", "_")).rstrip()
    return safe.strip().replace(" ", "_")

def extract_text(path: Path) -> str:
    ext = path.suffix.lower()
    if ext in ['.txt', '.md', '.boxnote']:
        return path.read_text(encoding="utf-8")
    if ext == '.xlsx':
        df = pd.read_excel(path)
        return df.to_csv(index=False)
    if ext == '.pdf':
        doc = fitz.open(str(path))
        return "".join(page.get_text() for page in doc)
    if ext == '.docx':
        docx_doc = docx.Document(str(path))
        return "\n".join(p.text for p in docx_doc.paragraphs)
    if ext == '.pptx':
        prs = Presentation(str(path))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    return ""

def build_vectorstore(notes_data, chunk_size: int = 1000):
    docs = []
    for note in notes_data:
        text = note["text"]
        # split into chunks
        for idx, i in enumerate(range(0, len(text), chunk_size)):
            chunk = text[i:i + chunk_size]
            docs.append(Document(page_content=chunk, metadata={"source": note["title"], "position": idx}))
    vs = FAISS.from_documents(docs, embeddings)
    vs.save_local(str(STORE_DIR))
    return vs

def load_vectorstore():
    if any(STORE_DIR.iterdir()):
        return FAISS.load_local(str(STORE_DIR), embeddings, allow_dangerous_deserialization=True)
    return None
    
st.set_page_config(page_title="NotebookLM風ノートアプリ", layout="wide")
st.title("NotebookLM風ノートアプリ (Streamlit + LangChain + FAISS)")

# Sidebar: note operations
st.sidebar.header("ノート操作")

# File upload
uploaded_files = st.sidebar.file_uploader(
    "ファイルをアップロード",
    type=['pdf', 'docx', 'txt', 'xlsx', 'pptx', 'boxnote'],
    accept_multiple_files=True
)
if uploaded_files:
    vs = load_vectorstore()
    for uploaded in uploaded_files:
        name = uploaded.name
        ext = name.split('.')[-1].lower()
        stem = name.rsplit(".", 1)[0]
        filename = slugify(stem) + "." + ext
        out_path = NOTES_DIR / filename
        with open(out_path, "wb") as f:
            f.write(uploaded.read())

        # Convert to markdown
        md_filename = slugify(stem) + ".md"
        md_path = MD_DIR / md_filename
        text = extract_text(out_path)
        md_path.write_text(text, encoding="utf-8")

        # Build documents for vectorstore
        docs = []
        chunk_size = 1000
        for idx, i in enumerate(range(0, len(text), chunk_size)):
            chunk = text[i:i + chunk_size]
            docs.append(Document(page_content=chunk, metadata={"source": str(out_path), "position": idx}))
        if vs:
            vs.add_documents(docs)
        else:
            vs = FAISS.from_documents(docs, embeddings)

    vs.save_local(str(STORE_DIR))
    st.sidebar.success("アップロードファイルをmarkdownに変換して保存し、ベクトルDBに登録しました")
    st.rerun()

# List notes
notes = sorted([f for f in NOTES_DIR.iterdir() if f.is_file()], key=lambda x: x.name)
note_names = [f.name for f in notes]
selected = st.sidebar.radio("ファイルを選択", [""] + note_names)

# New text note
new_title = st.sidebar.text_input("新規ノート作成", "")
if st.sidebar.button("作成"):
    if new_title.strip():
        filename = slugify(new_title) + ".txt"
        p = NOTES_DIR / filename
        if p.exists():
            st.sidebar.warning(f"ノート '{new_title}' は既に存在します。")
        else:
            p.write_text("", encoding="utf-8")
            for f in STORE_DIR.iterdir():
                f.unlink(missing_ok=True)
            st.sidebar.success(f"作成しました: {filename}")
            st.rerun()

# Delete note
if selected:
    st.sidebar.markdown("---")
    if st.sidebar.button("削除"):
        (NOTES_DIR / selected).unlink(missing_ok=True)
        for f in STORE_DIR.iterdir():
            f.unlink(missing_ok=True)
        st.sidebar.success(f"削除しました: {selected}")
        st.rerun()

# Main area: view/edit and RAG
if selected:
    path = NOTES_DIR / selected
    ext = path.suffix.lower()
    content = extract_text(path)
    if ext in ['.txt', '.md', '.boxnote']:
        st.subheader(f"編集: {selected}")
        edited = st.text_area("ファイル編集", value=content, height=300)
        if st.button("保存"):
            path.write_text(edited, encoding="utf-8")
            for f in STORE_DIR.iterdir():
                f.unlink(missing_ok=True)
            st.success("保存しました。")
    else:
        st.subheader(f"プレビュー: {selected}")
        st.text_area("内容", value=content, height=300)
else:
    st.info("左側のサイドバーからファイルを選択するか、新規作成してください。")
    st.markdown("---")
    st.subheader("AIに質問 (RAG)")
    system_prompt = st.text_area(
        "システムプロンプト",
        value="あなたはノートアプリのアシスタントです。",
        height=100
    )
    user_question = st.text_area(
        "質問",
        value="このノートの要約を教えてください",
        height=150
    )
    if st.button("質問する"):
        notes_data = [{"title": f.stem, "text": extract_text(f)} for f in notes]
        if not notes_data:
            st.warning("ノートが登録されていません。まずノートを作成またはアップロードしてください。")
        else:
            vs = load_vectorstore()
            st.write("vectorstoreを読み込み")
            if vs is None:
                vs = build_vectorstore(notes_data)
            docs = vs.similarity_search(user_question, k=7)
            st.write("類似テキストを読み込み",docs)
            context = ""
            for doc in docs:
                src = doc.metadata.get("source", "")
                context += f"### {src}\n{doc.page_content}\n\n"
            prompt = (
                "次のノート群を参照して、質問に答えてください。\n\n"
                f"{context}"
                f"### 質問:\n{user_question}\n"
            )
            resp = client.chat.completions.create(
                model=CHAT_DEPLOYMENT,
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7,
                max_tokens=512
            )
            answer = resp.choices[0].message.content.strip()
            st.subheader("AIの回答")
            st.markdown(answer)

            st.markdown("#### 参照ノート情報")
            for doc in docs:
                src = doc.metadata.get("source", "")
                st.markdown(f"- **{src}**")
                st.markdown(f"{doc.page_content}")
